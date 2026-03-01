"""
Office Integration Tools

Provides tools for working with Microsoft Office files:
- Word documents (.docx)
- Excel spreadsheets (.xlsx)
- PowerPoint presentations (.pptx)

Features:
- Read/write/modify documents
- Create new documents from templates
- Extract text and data
- Format text and cells
"""

import os
import logging
from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from enum import Enum
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """Document types"""
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    UNKNOWN = "unknown"


@dataclass
class DocumentResult:
    """Document operation result"""
    success: bool
    file_path: Optional[str] = None
    content: Optional[str] = None
    data: Optional[Any] = None
    error: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None


class WordTool:
    """Microsoft Word document tool"""
    
    def __init__(self):
        self._docx_available = self._check_docx()
    
    def _check_docx(self) -> bool:
        """Check if python-docx is available"""
        try:
            import docx
            return True
        except ImportError:
            logger.warning("python-docx not installed. Install with: pip install python-docx")
            return False
    
    async def create_document(
        self,
        file_path: str,
        title: str = "",
        paragraphs: List[str] = None
    ) -> DocumentResult:
        """Create a new Word document"""
        if not self._docx_available:
            return DocumentResult(
                success=False,
                error="python-docx not installed"
            )
        
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # Add title
            if title:
                heading = doc.add_heading(title, 0)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Add paragraphs
            if paragraphs:
                for para in paragraphs:
                    doc.add_paragraph(para)
            
            # Save
            doc.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                metadata={"paragraphs": len(paragraphs) if paragraphs else 0}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def read_document(self, file_path: str) -> DocumentResult:
        """Read text from Word document"""
        if not self._docx_available:
            return DocumentResult(
                success=False,
                error="python-docx not installed"
            )
        
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                content="\n\n".join(paragraphs),
                metadata={"paragraph_count": len(paragraphs)}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def add_paragraph(
        self,
        file_path: str,
        text: str,
        style: str = "Normal",
        bold: bool = False,
        italic: bool = False,
        font_size: int = 12
    ) -> DocumentResult:
        """Add a paragraph to existing document"""
        if not self._docx_available:
            return DocumentResult(
                success=False,
                error="python-docx not installed"
            )
        
        try:
            from docx import Document
            from docx.shared import Pt
            
            doc = Document(file_path)
            para = doc.add_paragraph(text, style=style)
            
            # Apply formatting
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            
            doc.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def create_table(
        self,
        file_path: str,
        table_data: List[List[str]],
        headers: List[str] = None
    ) -> DocumentResult:
        """Create a table in Word document"""
        if not self._docx_available:
            return DocumentResult(
                success=False,
                error="python-docx not installed"
            )
        
        try:
            from docx import Document
            
            doc = Document(file_path) if os.path.exists(file_path) else Document()
            
            # Add table
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            
            if headers and cols > 0:
                table = doc.add_table(rows=rows + 1, cols=cols)
                # Add header row
                header_row = table.rows[0]
                for i, header in enumerate(headers[:cols]):
                    header_row.cells[i].text = header
                # Add data rows
                for i, row_data in enumerate(table_data):
                    row = table.rows[i + 1]
                    for j, cell_data in enumerate(row_data[:cols]):
                        row.cells[j].text = str(cell_data)
            else:
                table = doc.add_table(rows=rows, cols=cols)
                for i, row_data in enumerate(table_data):
                    row = table.rows[i]
                    for j, cell_data in enumerate(row_data[:cols]):
                        row.cells[j].text = str(cell_data)
            
            doc.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                metadata={"rows": rows, "cols": cols}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))


class ExcelTool:
    """Microsoft Excel spreadsheet tool"""
    
    def __init__(self):
        self._openpyxl_available = self._check_openpyxl()
    
    def _check_openpyxl(self) -> bool:
        """Check if openpyxl is available"""
        try:
            import openpyxl
            return True
        except ImportError:
            logger.warning("openpyxl not installed. Install with: pip install openpyxl")
            return False
    
    async def create_workbook(
        self,
        file_path: str,
        sheet_name: str = "Sheet1",
        data: List[List[Any]] = None
    ) -> DocumentResult:
        """Create a new Excel workbook"""
        if not self._openpyxl_available:
            return DocumentResult(
                success=False,
                error="openpyxl not installed"
            )
        
        try:
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            # Add data
            if data:
                for row_idx, row_data in enumerate(data, start=1):
                    for col_idx, cell_data in enumerate(row_data, start=1):
                        ws.cell(row=row_idx, column=col_idx, value=cell_data)
            
            wb.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                metadata={"rows": len(data) if data else 0, "sheet": sheet_name}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def read_sheet(
        self,
        file_path: str,
        sheet_name: Optional[str] = None,
        as_dict: bool = False
    ) -> DocumentResult:
        """Read data from Excel sheet"""
        if not self._openpyxl_available:
            return DocumentResult(
                success=False,
                error="openpyxl not installed"
            )
        
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            ws = wb[sheet_name] if sheet_name else wb.active
            
            if as_dict:
                # Read as list of dictionaries
                headers = [cell.value for cell in ws[1]]
                rows = []
                for row in ws.iter_rows(min_row=2, values_only=True):
                    rows.append(dict(zip(headers, row)))
                data = rows
            else:
                # Read as list of lists
                data = [[cell.value for cell in row] for row in ws.iter_rows()]
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                data=data,
                metadata={
                    "sheet_name": ws.title,
                    "rows": ws.max_row,
                    "cols": ws.max_column
                }
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def write_cells(
        self,
        file_path: str,
        sheet_name: str,
        start_row: int,
        start_col: int,
        data: List[List[Any]]
    ) -> DocumentResult:
        """Write data to cells"""
        if not self._openpyxl_available:
            return DocumentResult(
                success=False,
                error="openpyxl not installed"
            )
        
        try:
            from openpyxl import load_workbook, Workbook
            
            if os.path.exists(file_path):
                wb = load_workbook(file_path)
            else:
                wb = Workbook()
            
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            
            ws = wb[sheet_name]
            
            for row_idx, row_data in enumerate(data, start=start_row):
                for col_idx, cell_data in enumerate(row_data, start=start_col):
                    ws.cell(row=row_idx, column=col_idx, value=cell_data)
            
            wb.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                metadata={"cells_written": len(data) * len(data[0]) if data else 0}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def create_chart(
        self,
        file_path: str,
        sheet_name: str,
        chart_type: str = "bar",
        title: str = "Chart",
        data_range: str = "A1:B10"
    ) -> DocumentResult:
        """Create a chart in Excel"""
        if not self._openpyxl_available:
            return DocumentResult(
                success=False,
                error="openpyxl not installed"
            )
        
        try:
            from openpyxl import load_workbook
            from openpyxl.chart import BarChart, LineChart, PieChart, Reference
            
            wb = load_workbook(file_path)
            ws = wb[sheet_name]
            
            # Parse range - simplified
            chart_types = {
                "bar": BarChart,
                "line": LineChart,
                "pie": PieChart
            }
            
            ChartClass = chart_types.get(chart_type, BarChart)
            chart = ChartClass()
            chart.title = title
            
            # Add data (simplified)
            data = Reference(ws, range_string=data_range)
            chart.add_data(data, titles_from_data=True)
            
            ws.add_chart(chart, "E5")
            
            wb.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))


class PowerPointTool:
    """Microsoft PowerPoint presentation tool"""
    
    def __init__(self):
        self._pptx_available = self._check_pptx()
    
    def _check_pptx(self) -> bool:
        """Check if python-pptx is available"""
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return False
    
    async def create_presentation(
        self,
        file_path: str,
        title: str = "",
        slides: List[Dict[str, Any]] = None
    ) -> DocumentResult:
        """Create a new PowerPoint presentation"""
        if not self._pptx_available:
            return DocumentResult(
                success=False,
                error="python-pptx not installed"
            )
        
        try:
            from pptx import Presentation
            
            prs = Presentation()
            
            # Set title slide
            if title:
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title_shape = slide.shapes.title
                subtitle = slide.placeholders[1]
                title_shape.text = title
                subtitle.text = "Created with RANGEN"
            
            # Add slides
            if slides:
                for slide_data in slides:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    if "title" in slide_data:
                        slide.shapes.title.text = slide_data["title"]
                    
                    if "content" in slide_data:
                        content = slide.placeholders[1]
                        text_frame = content.text_frame
                        text_frame.text = slide_data["content"]
            
            prs.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                metadata={"slides": len(prs.slides)}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def read_slides(self, file_path: str) -> DocumentResult:
        """Read text from PowerPoint slides"""
        if not self._pptx_available:
            return DocumentResult(
                success=False,
                error="python-pptx not installed"
            )
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text:
                            slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            
            return DocumentResult(
                success=True,
                file_path=file_path,
                content="\n\n---\n\n".join(slides_text),
                metadata={"slide_count": len(prs.slides)}
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))
    
    async def add_slide(
        self,
        file_path: str,
        title: str,
        content: str = "",
        bullet_points: List[str] = None
    ) -> DocumentResult:
        """Add a slide to presentation"""
        if not self._pptx_available:
            return DocumentResult(
                success=False,
                error="python-pptx not installed"
            )
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path) if os.path.exists(file_path) else Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.title.text = title
            
            if content or bullet_points:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                if content:
                    text_frame.text = content
                
                if bullet_points:
                    for point in bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 1
            
            prs.save(file_path)
            
            return DocumentResult(
                success=True,
                file_path=file_path
            )
            
        except Exception as e:
            return DocumentResult(success=False, error=str(e))


class OfficeTools:
    """
    Unified Office tools interface
    Provides a single interface for all Office operations
    """
    
    def __init__(self):
        self.word = WordTool()
        self.excel = ExcelTool()
        self.powerpoint = PowerPointTool()
    
    @staticmethod
    def detect_type(file_path: str) -> DocumentType:
        """Detect document type from extension"""
        ext = Path(file_path).suffix.lower()
        
        if ext == ".docx":
            return DocumentType.WORD
        elif ext == ".xlsx":
            return DocumentType.EXCEL
        elif ext == ".pptx":
            return DocumentType.POWERPOINT
        else:
            return DocumentType.UNKNOWN
    
    async def process(
        self,
        operation: str,
        file_path: str,
        **kwargs
    ) -> DocumentResult:
        """Process document based on detected type"""
        doc_type = self.detect_type(file_path)
        
        if doc_type == DocumentType.WORD:
            return await self._process_word(operation, file_path, **kwargs)
        elif doc_type == DocumentType.EXCEL:
            return await self._process_excel(operation, file_path, **kwargs)
        elif doc_type == DocumentType.POWERPOINT:
            return await self._process_powerpoint(operation, file_path, **kwargs)
        else:
            return DocumentResult(
                success=False,
                error=f"Unsupported document type: {doc_type}"
            )
    
    async def _process_word(
        self,
        operation: str,
        file_path: str,
        **kwargs
    ) -> DocumentResult:
        """Process Word document"""
        if operation == "create":
            return await self.word.create_document(file_path, **kwargs)
        elif operation == "read":
            return await self.word.read_document(file_path)
        elif operation == "add_paragraph":
            return await self.word.add_paragraph(file_path, **kwargs)
        elif operation == "add_table":
            return await self.word.create_table(file_path, **kwargs)
        else:
            return DocumentResult(success=False, error=f"Unknown operation: {operation}")
    
    async def _process_excel(
        self,
        operation: str,
        file_path: str,
        **kwargs
    ) -> DocumentResult:
        """Process Excel spreadsheet"""
        if operation == "create":
            return await self.excel.create_workbook(file_path, **kwargs)
        elif operation == "read":
            return await self.excel.read_sheet(file_path, **kwargs)
        elif operation == "write":
            return await self.excel.write_cells(file_path, **kwargs)
        elif operation == "chart":
            return await self.excel.create_chart(file_path, **kwargs)
        else:
            return DocumentResult(success=False, error=f"Unknown operation: {operation}")
    
    async def _process_powerpoint(
        self,
        operation: str,
        file_path: str,
        **kwargs
    ) -> DocumentResult:
        """Process PowerPoint presentation"""
        if operation == "create":
            return await self.powerpoint.create_presentation(file_path, **kwargs)
        elif operation == "read":
            return await self.powerpoint.read_slides(file_path)
        elif operation == "add_slide":
            return await self.powerpoint.add_slide(file_path, **kwargs)
        else:
            return DocumentResult(success=False, error=f"Unknown operation: {operation}")


# ==================== Convenience Functions ====================

_office_tools: Optional[OfficeTools] = None


def get_office_tools() -> OfficeTools:
    """Get Office tools instance"""
    global _office_tools
    if _office_tools is None:
        _office_tools = OfficeTools()
    return _office_tools
